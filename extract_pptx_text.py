import os
import sys
from pptx import Presentation

# Function to extract text from a slide
def extract_text_from_slide(slide):
    text = []
    for shape in slide.shapes:
        # Check if the shape has text, and append it to a list if it does
        if hasattr(shape, 'text'):
            text.append(shape.text)
    # Join the text from all shapes into a single string with newlines
    return '\n'.join(text)

# Function to extract notes from a slide
def extract_notes_from_slide(slide):
    # Get the notes slide for the given slide
    notes_slide = slide.notes_slide
    if notes_slide is not None:
        # If the notes slide exists, extract its text
        return extract_text_from_slide(notes_slide)
    # If there are no notes, return an empty string
    return ""

# Main function to extract text and notes from a PowerPoint file
def main(pptx_file, output_file):
    # Check if the PowerPoint file exists
    if not os.path.exists(pptx_file):
        print(f"Error: File '{pptx_file}' does not exist.")
        sys.exit(1)

    # Load the PowerPoint file as a presentation object
    presentation = Presentation(pptx_file)

    # Open the output file for writing
    with open(output_file, 'w', encoding='utf-8') as f:
        # Iterate over each slide in the presentation
        for slide_num, slide in enumerate(presentation.slides, start=1):
            # Extract the text and notes from the slide
            slide_text = extract_text_from_slide(slide)
            slide_notes = extract_notes_from_slide(slide)
            # Write the slide number, text, and notes to the output file
            f.write(f"Slide {slide_num}:\n{slide_text}\n\nNotes:\n{slide_notes}\n\n")
            # Print a message indicating the slide and its notes were extracted
            print(f"Slide {slide_num} and its notes extracted.")

    # Print a message indicating the extraction is complete
    print(f"Extraction completed. Results saved in '{output_file}'.")

# Check if the script was run as the main program
if __name__ == '__main__':
    # Check that the script was given the correct number of arguments
    if len(sys.argv) != 3:
        print("Usage: python extract_pptx_text.py <pptx_file> <output_file>")
        sys.exit(1)

    # Get the input and output files from the command line arguments
    pptx_file = sys.argv[1]
    output_file = sys.argv[2]
    # Call the main function with the input and output files
    main(pptx_file, output_file)
